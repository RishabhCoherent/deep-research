"""
Extract Table of Contents from a PPTX report file.

Discovery-based parser — uses paragraph margin-left (marL) from the PPTX
XML to determine entry hierarchy. No hardcoded region names, segment names,
year ranges, section titles, or font sizes.

Hierarchy determined by marL indentation values:
  - Level 0 (marL=0): Section headers, section titles, continuation slides
  - Level 1+ (marL>0): Subsections at increasing nesting depth

Section type classification (from title patterns):
  - Segment sections: title contains ", By <Dimension>,"
  - Region sections: segment dimension starts with "Region"
  - Competitive sections: title contains "competitive" or "company profile"
  - All others: subsection-based (nested tree)
"""

import re
import logging
from pathlib import Path

from pptx import Presentation

logger = logging.getLogger(__name__)


# ─── Patterns ────────────────────────────────────────────────────────────────

# Section header: "Section 1", "Chapter 2", "Part 10"
_SECTION_HEADER_RE = re.compile(r'^(?:section|chapter|part)\s+(\d+)$', re.I)

# Year range: "2020 - 2032", "2024–2030"
_YEAR_RANGE_RE = re.compile(r'\d{4}\s*[-–—]\s*\d{2,4}')

# Segment section title: "..., By Product Type, 2020 - 2032, ..."
_BY_DIMENSION_IN_TITLE_RE = re.compile(r',\s*by\s+([^,]+?)\s*,\s*\d{4}', re.I)

# Competitive section keywords
_COMPETITIVE_KEYWORDS = re.compile(
    r'competitive\s+landscape|competitive\s+analysis|company\s+profile', re.I
)

# Company group header: "Global Players", "North America Players"
_COMPANY_GROUP_RE = re.compile(r'^(.+?)\s+players$', re.I)

# TOC slide title keywords
_TOC_TITLE_KEYWORDS = re.compile(
    r'table\s+of\s+contents|^toc$|^contents$', re.I
)

# Continuation marker
_CONTINUATION_RE = re.compile(r'continu', re.I)

# Structural/skip entries
_SKIP_PATTERNS = [
    re.compile(r'^introduction$', re.I),
    re.compile(r'^segment\s+trends$', re.I),
    re.compile(r'^regional\s+trends$', re.I),
    re.compile(r'^market\s+share\s*\(%?\)\s*analysis', re.I),
    re.compile(r'^market\s+y-o-y\s+growth', re.I),
    re.compile(r'^market\s+size\s+and\s+forecast', re.I),
    re.compile(r'^company\s+overview$', re.I),
    re.compile(r'^product.{0,5}service\s+portfolio$', re.I),
    re.compile(r'^financial\s+performance$', re.I),
    re.compile(r'^recent\s+development', re.I),
    re.compile(r'^strategic\s+overview$', re.I),
    re.compile(r'^what\s+market\s+participants', re.I),
    re.compile(r'^competitive\s+dashboard$', re.I),
    re.compile(r'^company\s+.*?(market\s+share|pricing)', re.I),
]


# ─── Metadata Extraction ────────────────────────────────────────────────────


def _get_report_metadata(prs: Presentation) -> dict:
    """Extract report title and subtitle from the first slide."""
    metadata = {"report_title": "", "subtitle": ""}
    if not prs.slides:
        return metadata

    slide = prs.slides[0]
    text_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        try:
            is_placeholder = shape.placeholder_format is not None
        except (ValueError, AttributeError):
            is_placeholder = False
        font_size = 0
        for para in shape.text_frame.paragraphs:
            if para.runs:
                sz = para.runs[0].font.size
                if sz:
                    font_size = max(font_size, sz)
                    break
        text_shapes.append({
            "text": text, "top": shape.top or 0,
            "font_size": font_size, "is_placeholder": is_placeholder,
        })

    text_shapes.sort(key=lambda s: (not s["is_placeholder"], -s["font_size"], s["top"]))

    for ts in text_shapes:
        text = ts["text"]
        if re.search(r'©|copyright|all\s+rights\s+reserved', text, re.I):
            continue
        if text.isdigit():
            continue
        if not metadata["report_title"]:
            metadata["report_title"] = text
        elif not metadata["subtitle"]:
            metadata["subtitle"] = text
            break

    return metadata


# ─── TOC Entry Extraction (with hierarchy) ──────────────────────────────────


def _is_toc_slide(slide) -> bool:
    """Check if a slide is a TOC slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if _TOC_TITLE_KEYWORDS.search(para.text.strip()):
                return True
    return False


def _extract_toc_entries(prs: Presentation) -> list[dict]:
    """Extract TOC entries with hierarchy info from paragraph marL.

    Returns list of:
        {"text": str, "marL": int, "bold": bool, "font_size": float|None}
    """
    entries = []
    for slide in prs.slides:
        if not _is_toc_slide(slide):
            continue
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if not text:
                        continue
                    para = cell.text_frame.paragraphs[0]
                    pPr = para._pPr
                    marL = int(pPr.get('marL', 0)) if pPr is not None and pPr.get('marL') else 0
                    font_size = None
                    bold = False
                    if para.runs:
                        font_size = para.runs[0].font.size
                        bold = para.runs[0].font.bold or False
                    entries.append({
                        "text": text,
                        "marL": marL,
                        "bold": bold,
                        "font_size": font_size,
                    })
                    break
    return entries


def _normalize_levels(entries: list[dict]) -> list[dict]:
    """Assign a normalized 'level' (0, 1, 2, ...) based on unique marL values."""
    unique_margins = sorted(set(e["marL"] for e in entries))
    margin_to_level = {m: i for i, m in enumerate(unique_margins)}
    for e in entries:
        e["level"] = margin_to_level[e["marL"]]
    return entries


# ─── Helpers ─────────────────────────────────────────────────────────────────


def _is_skip_entry(text: str) -> bool:
    """Check if a TOC entry is structural boilerplate."""
    return any(p.search(text) for p in _SKIP_PATTERNS)


def _is_continuation(text: str) -> bool:
    """Check if a TOC entry is a continuation marker."""
    return bool(_CONTINUATION_RE.search(text) and '..' in text)


def _clean_year_range(text: str) -> str:
    """Remove year range: 'Marker Balls, 2020 - 2032' -> 'Marker Balls'."""
    return _YEAR_RANGE_RE.sub('', text).rstrip(' ,').strip()


def _build_nested(entries: list[tuple[str, int]]) -> list[dict]:
    """Build a nested tree from (text, level) tuples.

    Each entry becomes {"title": str} with optional "children": [...].
    Children are entries at a deeper level immediately following.
    """
    if not entries:
        return []

    result = []
    i = 0
    while i < len(entries):
        text, level = entries[i]
        # Collect children: all entries at deeper level until same/shallower level
        children_entries = []
        j = i + 1
        while j < len(entries) and entries[j][1] > level:
            children_entries.append(entries[j])
            j += 1

        node = {"title": text}
        if children_entries:
            node["children"] = _build_nested(children_entries)
        result.append(node)
        i = j

    return result


# ─── Section Parsing ────────────────────────────────────────────────────────


def _parse_sections(entries: list[dict]) -> list[dict]:
    """Parse entries into hierarchical sections using marL-based levels."""
    sections = []
    current_section = None
    current_entries = []

    def _flush_section():
        nonlocal current_section, current_entries
        if current_section is None:
            return

        title = current_section.get("title")

        if title and _COMPETITIVE_KEYWORDS.search(title):
            _parse_competitive_section(current_section, current_entries)
        elif title and _BY_DIMENSION_IN_TITLE_RE.search(title):
            dim_match = _BY_DIMENSION_IN_TITLE_RE.search(title)
            dim = dim_match.group(1).strip()
            current_section["segment_dimension"] = dim
            if re.match(r'^region', dim, re.I):
                _parse_region_section(current_section, current_entries)
            else:
                _parse_segment_section(current_section, current_entries)
        else:
            _parse_subsection_section(current_section, current_entries)

        sections.append(current_section)
        current_section = None
        current_entries = []

    for entry in entries:
        text = entry["text"]
        level = entry["level"]
        bold = entry["bold"]

        # Skip continuation slides (L0, not bold)
        if level == 0 and not bold:
            continue

        # Section header: L0, bold, matches "Section N"
        m = _SECTION_HEADER_RE.match(text)
        if m and level == 0 and bold:
            _flush_section()
            current_section = {"section_number": int(m.group(1))}
            current_entries = []
            continue

        if current_section is None:
            continue

        # Section title: L0, bold (first one after section header)
        if level == 0 and bold:
            if "title" not in current_section:
                current_section["title"] = text
            else:
                # Additional L0 bold entries (e.g., company group headers)
                current_entries.append(entry)
            continue

        # All other entries belong to the current section
        current_entries.append(entry)

    _flush_section()
    return sections


def _parse_subsection_section(section: dict, entries: list[dict]):
    """Parse a generic section with nested subsections."""
    filtered = []
    for e in entries:
        text = e["text"]
        if _is_skip_entry(text) or _is_continuation(text):
            continue
        if _YEAR_RANGE_RE.search(text):
            continue
        filtered.append((text, e["level"]))

    section["subsections"] = _build_nested(filtered)


def _parse_segment_section(section: dict, entries: list[dict]):
    """Parse a segment section — extract segment names from L2+ entries."""
    levels = [e["level"] for e in entries if e["level"] > 0]
    if not levels:
        section["segments"] = []
        return

    min_level = min(levels)
    segment_level = min_level + 1

    segments = []
    for e in entries:
        if e["level"] == segment_level:
            name = _clean_year_range(e["text"])
            if name and not _is_skip_entry(name):
                segments.append(name)

    seen = set()
    unique_segments = []
    for s in segments:
        if s not in seen:
            seen.add(s)
            unique_segments.append(s)

    section["segments"] = unique_segments


def _parse_region_section(section: dict, entries: list[dict]):
    """Parse a region section — regions at L2, countries at deepest level."""
    levels = [e["level"] for e in entries if e["level"] > 0]
    if not levels:
        section["regions"] = []
        return

    min_level = min(levels)
    region_level = min_level + 1

    all_levels = sorted(set(e["level"] for e in entries if e["level"] > region_level))
    country_level = all_levels[-1] if all_levels else None

    regions = []
    current_region = None

    for e in entries:
        text = e["text"]
        if _is_skip_entry(text) or _is_continuation(text) or _YEAR_RANGE_RE.search(text):
            continue

        if e["level"] == region_level:
            current_region = {"name": text, "countries": []}
            regions.append(current_region)
        elif country_level and e["level"] == country_level and current_region is not None:
            current_region["countries"].append(text)

    section["regions"] = regions


def _parse_competitive_section(section: dict, entries: list[dict]):
    """Parse competitive landscape — company groups and company names."""
    subsections = []
    companies = {}
    current_group = None

    for e in entries:
        text = e["text"]
        level = e["level"]
        bold = e["bold"]

        if _is_continuation(text):
            continue

        # Group header: L0 bold "... Players"
        if level == 0 and bold:
            gm = _COMPANY_GROUP_RE.match(text)
            if gm:
                current_group = re.sub(r'[^a-z0-9]+', '_', gm.group(1).lower()).strip('_')
                companies.setdefault(current_group, [])
            continue

        # Before any company group → analytical subsections
        if current_group is None:
            if _is_skip_entry(text):
                continue
            if level >= 1:
                subsections.append({"title": text})
            continue

        # Company names at L1 (skip profile items at L2+)
        if level == 1:
            company = text.rstrip("*").strip()
            if company and not _is_skip_entry(company):
                companies[current_group].append(company)

    if subsections:
        section["subsections"] = subsections
    section["companies"] = companies


# ─── Public API ──────────────────────────────────────────────────────────────


def extract_toc(pptx_path: str | Path) -> dict:
    """Extract the Table of Contents from a PPTX report file.

    Uses paragraph margin-left (marL) to determine hierarchy. Section types
    (segment, region, competitive) are classified from title patterns.

    Args:
        pptx_path: Path to the PPTX file.

    Returns:
        Dictionary with report metadata and hierarchical TOC structure.
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    logger.info(f"Extracting TOC from: {pptx_path.name}")

    prs = Presentation(str(pptx_path))

    metadata = _get_report_metadata(prs)

    entries = _extract_toc_entries(prs)
    logger.info(f"Found {len(entries)} raw TOC entries")

    entries = _normalize_levels(entries)

    sections = _parse_sections(entries)
    logger.info(f"Parsed into {len(sections)} sections")

    return {
        "report_title": metadata["report_title"],
        "subtitle": metadata["subtitle"],
        "total_slides": len(prs.slides),
        "sections": sections,
    }
